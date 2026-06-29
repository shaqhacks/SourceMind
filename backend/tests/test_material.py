"""Unit tests for backend/extract/material.py — TDD: tests written first."""
from __future__ import annotations

import pytest
from pathlib import Path


# ---------------------------------------------------------------------------
# paginate_text
# ---------------------------------------------------------------------------

class TestPaginateText:
    def test_paginate_text_basic(self):
        from SourceMind.backend.extract.material import paginate_text

        words = list(range(10))
        text = " ".join(str(w) for w in words)  # 10 words
        pages = paginate_text(text, words_per_page=3)

        assert len(pages) == 4  # 3/3/3/1
        # sequential page numbers are not part of paginate_text (it returns str),
        # but let's verify chunks
        assert pages[0] == "0 1 2"
        assert pages[1] == "3 4 5"
        assert pages[2] == "6 7 8"
        assert pages[3] == "9"

    def test_paginate_text_exact_budget(self):
        from SourceMind.backend.extract.material import paginate_text

        n, budget = 4, 5
        text = " ".join(["word"] * (n * budget))
        pages = paginate_text(text, words_per_page=budget)

        assert len(pages) == n

    def test_paginate_text_empties_dropped(self):
        from SourceMind.backend.extract.material import paginate_text

        # Extra whitespace / blank lines should produce no empty chunks
        text = "  hello   world  \n\n\n  foo  \n"
        pages = paginate_text(text, words_per_page=10)

        assert all(p.strip() for p in pages)

    def test_paginate_text_no_mid_word_splits(self):
        from SourceMind.backend.extract.material import paginate_text

        original_words = [f"token{i}" for i in range(50)]
        text = " ".join(original_words)
        pages = paginate_text(text, words_per_page=7)

        all_words_in_pages = set()
        for p in pages:
            all_words_in_pages.update(p.split())

        assert all_words_in_pages == set(original_words)


# ---------------------------------------------------------------------------
# detect_kind
# ---------------------------------------------------------------------------

@pytest.mark.parametrize(
    "filename,expected_kind",
    [
        ("lecture.pdf", "pdf"),
        ("notes.docx", "docx"),
        ("slides.pptx", "pptx"),
        ("readme.txt", "txt"),
        ("chapter.md", "md"),
    ],
)
def test_detect_kind_extensions(filename, expected_kind):
    from SourceMind.backend.extract.material import detect_kind

    assert detect_kind(filename) == expected_kind


def test_detect_kind_unknown_raises():
    from SourceMind.backend.extract.material import detect_kind

    with pytest.raises(ValueError, match="Unsupported extension"):
        detect_kind("file.xyz")


# ---------------------------------------------------------------------------
# extract_docx
# ---------------------------------------------------------------------------

def test_extract_docx(tmp_path):
    import docx as python_docx
    from SourceMind.backend.extract.material import extract_docx

    doc = python_docx.Document()
    doc.add_paragraph("Hello from docx paragraph one.")
    doc.add_paragraph("Second paragraph here.")
    dest = tmp_path / "test.docx"
    doc.save(str(dest))

    result = extract_docx(dest)

    assert "Hello from docx paragraph one." in result
    assert "Second paragraph here." in result


# ---------------------------------------------------------------------------
# extract_pptx
# ---------------------------------------------------------------------------

def test_extract_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from SourceMind.backend.extract.material import extract_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txBox.text_frame
    tf.text = "Hello from pptx slide."
    dest = tmp_path / "test.pptx"
    prs.save(str(dest))

    result = extract_pptx(dest)

    assert "Hello from pptx slide." in result


def test_extract_pptx_multi_run_paragraph(tmp_path):
    """A paragraph with two runs must not gain a spurious newline between them."""
    from pptx import Presentation
    from pptx.util import Inches
    from SourceMind.backend.extract.material import extract_pptx

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = txBox.text_frame
    # Use the initial paragraph and add two runs to it
    para = tf.paragraphs[0]
    run1 = para.add_run()
    run1.text = "Hello, "
    run2 = para.add_run()
    run2.text = "world!"
    dest = tmp_path / "multi_run.pptx"
    prs.save(str(dest))

    result = extract_pptx(dest)

    assert "Hello, world!" in result
    assert "Hello, \nworld!" not in result


# ---------------------------------------------------------------------------
# extract_material — txt
# ---------------------------------------------------------------------------

def test_extract_material_txt_paginates(tmp_path):
    from SourceMind.backend.extract.material import extract_material

    # ~1500 words → 3 pages at default 500 words/page
    content = ("alpha beta gamma delta epsilon " * 300).strip()
    txt_file = tmp_path / "lecture.txt"
    txt_file.write_text(content)

    pages = extract_material("txt", path=txt_file)

    assert len(pages) >= 2
    # sequential page numbers starting at 0
    for i, page in enumerate(pages):
        assert page.page_number == i
    # all text recovered
    recovered = " ".join(p.text for p in pages)
    assert "alpha" in recovered and "epsilon" in recovered
    # no images for txt
    for page in pages:
        assert page.image_paths == []


# ---------------------------------------------------------------------------
# extract_material — pdf (delegates to extract_pdf)
# ---------------------------------------------------------------------------

def test_extract_material_pdf_delegates(tmp_path):
    import fitz
    from SourceMind.backend.extract.material import extract_material
    from SourceMind.backend.extract.pdf import ExtractedPage

    # Build a minimal single-page PDF with fitz
    pdf_path = tmp_path / "test.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 72), "PDF content for test.")
    doc.save(str(pdf_path))
    doc.close()

    assets_dir = tmp_path / "assets"
    pages = extract_material("pdf", path=pdf_path, assets_dir=assets_dir)

    assert len(pages) >= 1
    assert all(isinstance(p, ExtractedPage) for p in pages)
    # image_paths is a list (may be empty for a simple page)
    assert all(isinstance(p.image_paths, list) for p in pages)
    assert "PDF content for test." in pages[0].text


# ---------------------------------------------------------------------------
# extract_material — text kind
# ---------------------------------------------------------------------------

def test_extract_material_text_kind():
    from SourceMind.backend.extract.material import extract_material

    pages = extract_material("text", text="word " * 600)

    assert len(pages) > 1
    for i, page in enumerate(pages):
        assert page.page_number == i
    for page in pages:
        assert page.image_paths == []
