"""Unified material extraction: PDF, DOCX, PPTX, TXT, MD, URL, text, YouTube.

All sources normalise to ``list[ExtractedPage]`` so the existing
outline/plan/chapter/RAG pipeline continues unchanged.
"""
from __future__ import annotations

import socket
from pathlib import Path
from urllib.parse import urljoin, urlparse

from SourceMind.backend.extract.pdf import ExtractedPage, derive_title_from_pdf, extract_pdf
from SourceMind.backend.services.ingest import SsrfError, normalize_source, validate_public_url

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

_EXT_TO_KIND: dict[str, str] = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".txt": "txt",
    ".md": "md",
}


# ---------------------------------------------------------------------------
# Text pagination
# ---------------------------------------------------------------------------

def paginate_text(text: str, words_per_page: int = 500) -> list[str]:
    """Split *text* into ~word-budget chunks on whitespace; drop empty chunks."""
    words = text.split()
    chunks: list[str] = []
    for i in range(0, len(words), words_per_page):
        chunk = " ".join(words[i : i + words_per_page])
        if chunk:
            chunks.append(chunk)
    return chunks


# ---------------------------------------------------------------------------
# Kind detection
# ---------------------------------------------------------------------------

def detect_kind(filename: str) -> str:
    """Map filename extension to kind string.

    Raises :class:`ValueError` for unknown extensions.
    """
    ext = Path(filename).suffix.lower()
    if ext not in _EXT_TO_KIND:
        raise ValueError(
            f"Unsupported extension: {ext!r}. Supported: {', '.join(sorted(_EXT_TO_KIND))}"
        )
    return _EXT_TO_KIND[ext]


# ---------------------------------------------------------------------------
# Format-specific raw extractors
# ---------------------------------------------------------------------------

def extract_docx(path) -> str:
    """Extract all paragraph text from a .docx file using python-docx."""
    import docx as python_docx  # lazy import so the dep is optional at import time

    doc = python_docx.Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pptx(path) -> str:
    """Extract all shape text from a .pptx file using python-pptx."""
    from pptx import Presentation  # lazy import

    prs = Presentation(str(path))
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text:
                        texts.append(para.text)
    return "\n".join(texts)


# ---------------------------------------------------------------------------
# Title derivation
# ---------------------------------------------------------------------------

def material_title(kind: str, *, path=None, url: str | None = None) -> str:
    """Derive a human-readable title for a material.

    - pdf: embedded metadata → prettified stem
    - docx/pptx/txt/md: path.stem with underscores/hyphens replaced by spaces
    - url: domain extracted from URL string
    - youtube: "YouTube video"
    - text: "Pasted text"
    """
    kind = kind.lower()
    if kind == "pdf":
        p = Path(path)
        return derive_title_from_pdf(p, p.name)
    if kind in ("docx", "pptx", "txt", "md", "markdown"):
        return Path(path).stem.replace("_", " ").replace("-", " ").strip()
    if kind == "url":
        try:
            return urlparse(url).netloc or url or "Web page"
        except Exception:
            return url or "Web page"
    if kind == "youtube":
        return "YouTube video"
    if kind in ("text", "pasted"):
        return "Pasted text"
    return "Untitled"


# ---------------------------------------------------------------------------
# SSRF-safe URL fetcher
# ---------------------------------------------------------------------------


def _resolve(host: str) -> list[str]:
    """Resolve *host* to unique IP address strings (IPv4+IPv6) for SSRF validation."""
    infos = socket.getaddrinfo(host, None)
    return list({info[4][0] for info in infos})


def _fetch_url_safely(url: str, *, max_bytes: int = 5_000_000) -> str:
    """Fetch *url* safely, enforcing SSRF guard, redirect re-validation, and size cap.

    1. Validates the initial URL against the SSRF guard (with DNS resolution).
    2. Fetches with redirects DISABLED; follows up to 5 hops manually,
       re-running the SSRF guard on each Location header.
    3. Raises ``ValueError`` if the redirect limit is exceeded or the body
       exceeds *max_bytes*.
    4. Returns the final response as text.
    """
    import httpx  # lazy import — not available in all test envs

    validate_public_url(url, resolver=_resolve)

    max_redirects = 5
    current_url = url
    resp = None

    with httpx.Client(follow_redirects=False, timeout=30) as client:
        for hop in range(max_redirects + 1):
            resp = client.get(current_url)
            if not resp.is_redirect:
                break
            location = resp.headers.get("location", "")
            if not location:
                break
            next_url = urljoin(current_url, location)
            validate_public_url(next_url, resolver=_resolve)
            current_url = next_url
            if hop == max_redirects:
                raise ValueError(
                    f"Too many redirects fetching {url!r} (limit: {max_redirects})"
                )

    if resp is None:
        raise ValueError(f"No response received from {url!r}")
    if len(resp.content) > max_bytes:
        raise ValueError(
            f"Response body too large ({len(resp.content)} bytes > {max_bytes} limit)"
        )
    return resp.text


# ---------------------------------------------------------------------------
# Unified extractor
# ---------------------------------------------------------------------------

def extract_material(
    kind: str,
    *,
    path=None,
    text: str | None = None,
    url: str | None = None,
    assets_dir=None,
) -> list[ExtractedPage]:
    """Unified extraction to ``list[ExtractedPage]``.

    Args:
        kind: One of ``pdf``, ``docx``, ``pptx``, ``txt``, ``md``,
              ``url``, ``youtube``, ``text``.
        path: File path (required for file-based kinds).
        text: Raw text (required for ``text`` kind).
        url: URL string (required for ``url`` / ``youtube`` kinds).
        assets_dir: Directory for extracted images (``pdf`` only).

    Returns:
        Ordered list of :class:`ExtractedPage` objects, ``page_number``
        starting at 0, ``image_paths`` empty for non-PDF sources.
    """
    kind = (kind or "").strip().lower()

    # --- PDF: real pages + images via PyMuPDF ---
    if kind == "pdf":
        if assets_dir is None:
            assets_dir = Path(path).parent / "assets"
        return extract_pdf(Path(path), Path(assets_dir))

    # --- DOCX ---
    if kind == "docx":
        raw = extract_docx(path)
        return _pages_from_text(raw)

    # --- PPTX ---
    if kind == "pptx":
        raw = extract_pptx(path)
        return _pages_from_text(raw)

    # --- TXT ---
    if kind == "txt":
        raw = Path(path).read_text(errors="replace")
        return _pages_from_text(raw)

    # --- Markdown ---
    if kind in ("md", "markdown"):
        raw = Path(path).read_text(errors="replace")
        return _pages_from_text(raw)

    # --- URL: fetch then normalise (SSRF-safe) ---
    if kind == "url":
        html = _fetch_url_safely(url)
        normalized = normalize_source("url", html)
        return _pages_from_text(normalized)

    # --- YouTube ---
    if kind == "youtube":
        raise ValueError(
            "YouTube transcript ingestion is not yet supported. "
            "Please paste the video transcript using the 'Paste text' option instead."
        )

    # --- Pasted text ---
    if kind in ("text", "pasted"):
        normalized = normalize_source("text", text)
        return _pages_from_text(normalized)

    raise ValueError(
        f"Unsupported kind: {kind!r}. "
        f"Supported: pdf, docx, pptx, txt, md, url, youtube, text."
    )


# ---------------------------------------------------------------------------
# Internal helper
# ---------------------------------------------------------------------------

def _pages_from_text(raw: str) -> list[ExtractedPage]:
    """Paginate *raw* text and wrap each chunk in an :class:`ExtractedPage`."""
    chunks = paginate_text(raw)
    return [
        ExtractedPage(page_number=i, text=chunk, image_paths=[])
        for i, chunk in enumerate(chunks)
    ]
